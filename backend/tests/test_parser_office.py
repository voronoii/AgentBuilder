from __future__ import annotations

from pathlib import Path

import pytest

from app.services.knowledge.parsers.csv_parser import CsvParser
from app.services.knowledge.parsers.pptx import PptxParser
from app.services.knowledge.parsers.xlsx import XlsxParser

FIX = Path(__file__).parent / "fixtures"


@pytest.mark.asyncio
async def test_csv_parser_reads_rows() -> None:
    parsed = await CsvParser().parse(FIX / "sample.csv")
    assert "alice" in parsed.text and "bob" in parsed.text
    assert parsed.metadata["row_count"] == 2


@pytest.mark.asyncio
async def test_xlsx_parser(tmp_path: Path) -> None:
    from openpyxl import Workbook

    p = tmp_path / "s.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["h1", "h2"])
    ws.append(["r1c1", "r1c2"])
    wb.save(p)
    parsed = await XlsxParser().parse(p)
    assert "r1c1" in parsed.text


@pytest.mark.asyncio
async def test_pptx_parser(tmp_path: Path) -> None:
    from pptx import Presentation

    p = tmp_path / "s.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 100, 100)
    tb.text_frame.text = "hello slide"
    prs.save(p)
    parsed = await PptxParser().parse(p)
    assert "hello slide" in parsed.text
