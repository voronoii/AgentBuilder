from __future__ import annotations

import asyncio
from pathlib import Path

from pptx import Presentation

from app.services.knowledge.parsers.base import ParsedDocument


class PptxParser:
    async def parse(self, path: Path) -> ParsedDocument:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._parse_sync, path)

    def _parse_sync(self, path: Path) -> ParsedDocument:
        prs = Presentation(str(path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"# slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in p.runs)
                        if text.strip():
                            lines.append(text)
        return ParsedDocument(
            text="\n".join(lines),
            metadata={"slide_count": len(prs.slides), "filename": path.name},
        )
